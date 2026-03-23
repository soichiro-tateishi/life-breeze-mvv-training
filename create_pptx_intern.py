# -*- coding: utf-8 -*-
"""
Life Breeze MVV研修（インターン向け）PowerPoint生成スクリプト
デザイン：既存PDF準拠（ネイビー・レッド・ホワイト / Noto Sans JP）
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ====================
# Life Breeze カラーパレット
# ====================
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BG = RGBColor(0xF5, 0xF7, 0xFA)              # 極薄グレー背景
C_NAVY = RGBColor(0x1B, 0x3A, 0x5C)             # ネイビー（メイン）
C_NAVY_DARK = RGBColor(0x0F, 0x25, 0x3C)        # ダークネイビー
C_NAVY_LIGHT = RGBColor(0xE8, 0xEF, 0xF5)       # 淡ネイビー（カード背景）
C_RED = RGBColor(0xC4, 0x1E, 0x3A)              # レッド（アクセント）
C_RED_LIGHT = RGBColor(0xFD, 0xE8, 0xEB)        # 淡レッド
C_BLUE = RGBColor(0x1A, 0x56, 0xDB)             # ブルー
C_BLUE_LIGHT = RGBColor(0xE8, 0xF0, 0xFE)       # 淡ブルー
C_GREEN = RGBColor(0x05, 0x96, 0x69)            # グリーン
C_GREEN_LIGHT = RGBColor(0xEC, 0xFD, 0xF5)      # 淡グリーン
C_AMBER = RGBColor(0xD9, 0x73, 0x06)            # アンバー
C_GRAY = RGBColor(0x6B, 0x72, 0x80)             # グレー
C_GRAY_LIGHT = RGBColor(0xF3, 0xF4, 0xF6)       # 淡グレー
C_TITLE = RGBColor(0x1A, 0x20, 0x2C)            # タイトルテキスト
C_BODY = RGBColor(0x37, 0x41, 0x51)             # 本文テキスト
C_SUB = RGBColor(0x6B, 0x72, 0x80)              # サブテキスト
C_BORDER = RGBColor(0xD1, 0xD5, 0xDB)           # 境界線

# フォント
FONT = "Noto Sans JP"
FONT_EN = "Inter"

# レイアウト定数 (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
ML = Inches(0.75)       # 左マージン
MR = Inches(0.75)       # 右マージン
CONTENT_W = Inches(11.833)  # コンテンツ幅
TOTAL_PAGES = 18

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ====================
# ヘルパー関数
# ====================
def set_bg(slide, color=C_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=12,
             bold=False, color=C_BODY, align=PP_ALIGN.LEFT,
             font=FONT, anchor=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return txBox


def add_para(tf, text, size=12, bold=False, color=C_BODY,
             align=PP_ALIGN.LEFT, font=FONT, space_before=None):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    if space_before:
        p.space_before = space_before
    return p


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def footer(slide, num):
    # 左: Life Breeze Inc.
    add_text(slide, Inches(0.5), Inches(7.05), Inches(3), Inches(0.35),
             "Life Breeze Inc.", size=8, color=C_SUB, font=FONT)
    # 右: ページ番号
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.35),
             f"{num}/{TOTAL_PAGES}", size=8, color=C_SUB,
             align=PP_ALIGN.RIGHT, font=FONT)


def logo_badge(slide):
    """右上のLife Breezeロゴバッジ"""
    # 赤い四角
    sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(11.6), Inches(0.35), Inches(0.22), Inches(0.22))
    sq.fill.solid()
    sq.fill.fore_color.rgb = C_RED
    sq.line.fill.background()
    # テキスト
    add_text(slide, Inches(11.9), Inches(0.3), Inches(1.2), Inches(0.3),
             "Life Breeze", size=9, bold=True, color=C_NAVY, font=FONT)


def title_bar(slide, title, subtitle_en=None):
    """赤アクセントライン付きタイトル"""
    # タイトルテキスト
    add_text(slide, ML, Inches(0.4), CONTENT_W, Inches(0.6),
             title, size=24, bold=True, color=C_TITLE, font=FONT)
    # 赤アクセントライン
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   ML, Inches(0.95), Inches(1.2), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = C_RED
    line.line.fill.background()
    if subtitle_en:
        add_text(slide, ML + Inches(1.4), Inches(0.88), Inches(6), Inches(0.3),
                 subtitle_en, size=10, color=C_SUB, font=FONT_EN)


def key_message(slide, text, y=Inches(1.2), size=20):
    """キーメッセージ（青い左バー付き）"""
    # 青い左バー
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  ML, y, Inches(0.08), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_NAVY
    bar.line.fill.background()
    # テキスト
    add_text(slide, ML + Inches(0.2), y, CONTENT_W - Inches(0.3), Inches(0.7),
             text, size=size, bold=True, color=C_TITLE, font=FONT)


def bottom_bar(slide, text, text_en=None, y=Inches(5.8)):
    """青背景のボトムメッセージバー"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), y, SLIDE_W, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(1.0)
    tf.margin_right = Inches(1.0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    if text_en:
        p2 = tf.add_paragraph()
        p2.text = text_en
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(0xA0, 0xC4, 0xE8)
        p2.font.name = FONT_EN
        p2.alignment = PP_ALIGN.CENTER


def info_card(slide, x, y, w, h, title_text, body_text, accent_color=C_NAVY,
              body_en=None, icon_text=None):
    """情報カード（白背景・左サイド色付き）"""
    # カード背景
    card_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card_bg.fill.solid()
    card_bg.fill.fore_color.rgb = C_WHITE
    card_bg.line.color.rgb = C_BORDER
    card_bg.line.width = Pt(0.5)
    # 左サイドバー
    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   x, y, Inches(0.06), h)
    side.fill.solid()
    side.fill.fore_color.rgb = accent_color
    side.line.fill.background()
    # アイコンサークル
    if icon_text:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         x + Inches(0.2), y + Inches(0.15),
                                         Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(
            min(accent_color.red + 0x30, 0xFF),
            min(accent_color.green + 0x30, 0xFF),
            min(accent_color.blue + 0x30, 0xFF)
        )
        circle.line.fill.background()
    # タイトル
    title_x = x + Inches(0.25) if not icon_text else x + Inches(0.75)
    add_text(slide, title_x, y + Inches(0.12), w - Inches(0.5), Inches(0.35),
             title_text, size=13, bold=True, color=C_TITLE, font=FONT)
    # ボディ
    body_y = y + Inches(0.45)
    txBox = add_text(slide, title_x, body_y, w - Inches(0.5), h - Inches(0.6),
                     body_text, size=10, color=C_BODY, font=FONT)
    if body_en:
        add_para(txBox.text_frame, body_en, size=9, color=C_SUB, font=FONT_EN)
    return card_bg


def step_card(slide, x, y, w, num, label, desc, accent_color, desc_en=None):
    """ステップカード（番号付き）"""
    h = Inches(0.9)
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_WHITE
    bg.line.color.rgb = C_BORDER
    bg.line.width = Pt(0.5)
    # 左色帯
    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), h)
    side.fill.solid()
    side.fill.fore_color.rgb = accent_color
    side.line.fill.background()
    # 番号サークル
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x + Inches(0.2), y + Inches(0.2),
                                     Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_color
    circle.line.fill.background()
    tf_c = circle.text_frame
    tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_c = tf_c.paragraphs[0]
    p_c.text = str(num)
    p_c.font.size = Pt(16)
    p_c.font.bold = True
    p_c.font.color.rgb = C_WHITE
    p_c.font.name = FONT
    p_c.alignment = PP_ALIGN.CENTER
    # ラベル
    add_text(slide, x + Inches(0.85), y + Inches(0.12), w - Inches(1.1), Inches(0.35),
             label, size=13, bold=True, color=C_TITLE, font=FONT)
    # 説明
    txBox = add_text(slide, x + Inches(0.85), y + Inches(0.45), w - Inches(1.1), Inches(0.4),
                     desc, size=9, color=C_SUB, font=FONT)
    if desc_en:
        add_para(txBox.text_frame, desc_en, size=8, color=C_GRAY, font=FONT_EN)


def workshop_header(slide, ws_num, title, color, time_min):
    """ワークショップスライドのヘッダー"""
    # 左サイドバー（全高）
    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(0), Inches(0.12), SLIDE_H)
    side.fill.solid()
    side.fill.fore_color.rgb = color
    side.line.fill.background()
    # WORKSHOPバッジ
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.3), Inches(0.35), Inches(1.6), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf_b = badge.text_frame
    tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_b.margin_left = Inches(0.1)
    p_b = tf_b.paragraphs[0]
    p_b.text = f"WORKSHOP {ws_num:02d}"
    p_b.font.size = Pt(10)
    p_b.font.bold = True
    p_b.font.color.rgb = C_WHITE
    p_b.font.name = FONT_EN
    p_b.alignment = PP_ALIGN.CENTER
    # タイトル
    add_text(slide, Inches(0.3), Inches(0.8), Inches(9), Inches(0.6),
             title, size=24, bold=True, color=C_TITLE, font=FONT)
    # 赤アクセントライン
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.3), Inches(1.35), Inches(1.2), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = C_RED
    line.line.fill.background()
    # タイムバッジ
    time_badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(11.0), Inches(0.35), Inches(1.8), Inches(0.45))
    time_badge.fill.solid()
    time_badge.fill.fore_color.rgb = C_NAVY_DARK
    time_badge.line.fill.background()
    tf_t = time_badge.text_frame
    tf_t.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_t = tf_t.paragraphs[0]
    p_t.text = f"  {time_min} min"
    p_t.font.size = Pt(14)
    p_t.font.bold = True
    p_t.font.color.rgb = C_WHITE
    p_t.font.name = FONT_EN
    p_t.alignment = PP_ALIGN.CENTER


# ====================
# Slide 01: 表紙
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_NAVY)
# 日付
add_text(s, ML, Inches(1.0), Inches(8), Inches(0.4),
         "2026.XX.XX  |  Intern Training", size=12, color=RGBColor(0xA0, 0xC4, 0xE8), font=FONT_EN)
# メインタイトル
add_text(s, ML, Inches(1.6), Inches(8), Inches(1.2),
         "Life Breeze\nMVV研修", size=48, bold=True, color=C_WHITE, font=FONT)
# 赤ライン
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                           ML, Inches(3.4), Inches(1.5), Inches(0.07))
line.fill.solid()
line.fill.fore_color.rgb = C_RED
line.line.fill.background()
# サブタイトル
add_text(s, ML, Inches(3.7), Inches(8), Inches(0.8),
         "MVVを知り、自分の行動につなげる\nLearn our MVV. Connect it to your actions.",
         size=16, color=RGBColor(0xCC, 0xDD, 0xEE), font=FONT)
# 右下ロゴ
sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                         Inches(10.5), Inches(6.0), Inches(0.28), Inches(0.28))
sq.fill.solid()
sq.fill.fore_color.rgb = C_RED
sq.line.fill.background()
add_text(s, Inches(10.85), Inches(5.95), Inches(2.2), Inches(0.4),
         "Life Breeze Inc.", size=12, bold=True, color=C_WHITE, font=FONT)
add_text(s, Inches(10.85), Inches(6.3), Inches(2.2), Inches(0.3),
         "Global Expansion Strategy Team", size=8, color=RGBColor(0x88, 0xAA, 0xCC), font=FONT_EN)
# 装飾サークル（右上）
circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Inches(9.5), Inches(0.8), Inches(3.5), Inches(3.5))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x22, 0x48, 0x6E)
circle.line.fill.background()
add_notes(s, "今日は60分でLife BreezeのMVVを学びます。暗記テストはしません。自分の言葉で語れるようになることがゴールです。（1分）")


# ====================
# Slide 02: 今日のゴール
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_WHITE)
logo_badge(s)
title_bar(s, "今日のゴール", "Today's Goal")
key_message(s, 'MVVを"暗記"ではなく\n"自分の行動指針"にする')

steps = [
    ("なぜこの事業をするのか", "起点", "Why this business?"),
    ("何を目指すのか", "MVV", "What we aim for"),
    ("どう動くのか", "行動への翻訳", "How we act"),
    ("自分のアクションを決める", "", "Decide your action"),
]
for i, (label, sub, en) in enumerate(steps):
    y = Inches(2.2) + Inches(1.15) * i
    accent = C_RED if i == 3 else C_NAVY
    desc = f"（{sub}）" if sub else ""
    step_card(s, ML, y, CONTENT_W, i + 1, label, desc, accent, en)

footer(s, 2)
add_notes(s, "最後に、自分の30日アクションを1つ決めて帰ってもらいます。ワークが3回あります。（2分）")


# ====================
# Slide 03: Life Breezeの原点
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_WHITE)
logo_badge(s)
title_bar(s, "Life Breezeの原点", "Our Origin")
key_message(s, "週8時間、手で洗濯する人がいる。", y=Inches(1.15), size=22)

cw = Inches(5.6)
# 左カード
info_card(s, ML, Inches(2.2), cw, Inches(2.5),
          "時間の喪失 / Time Lost",
          "洗濯に時間が取られ、\n学び・仕事・家族・夢が後回しになる",
          C_NAVY,
          "Children can't study. Mothers can't work.\nDreams are put on hold.")
# 右カード
info_card(s, ML + Inches(6.0), Inches(2.2), cw, Inches(2.5),
          "インフラの不安定さ / Unstable Infrastructure",
          "電気・水の供給が不安定で、\n家事負担がさらに増大している",
          C_AMBER)

bottom_bar(s, '洗濯は"家事"ではなく、挑戦の時間を奪う「インフラ課題」だった。',
           "Laundry was not a chore — it was an infrastructure problem stealing time for ambition.",
           y=Inches(5.4))
footer(s, 3)
add_notes(s, "週8時間という数字をゆっくり伝える。「みなさんの自由時間は1日何時間？」と問いかけ、想像させる間を取る。同情ではなく課題として客観的に。（4分）")


# ====================
# Slide 04: 私たちの事業は何か
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_WHITE)
logo_badge(s)
title_bar(s, "私たちの事業は何か", "What We Do")
key_message(s, '"店舗を増やす"のではなく\n"自走できる仕組み"を届ける', y=Inches(1.15))

cards = [
    ("手段と目的 / Means & Purpose", "店舗は「手段」。\n目的は、現地の人の生活の自由度を上げること。", C_NAVY),
    ("自走する仕組み / Self-reliance", "日本人がいなくても回る状態を\n作ることが成功。", C_RED),
    ("展開状況 / Operations", "稼働中: モザンビーク、ケニア（4店舗）\n準備中: リベリア、ガボン、マダガスカル\n対象人口: 約8,700万人", C_BLUE),
]
for i, (title, body, color) in enumerate(cards):
    info_card(s, ML, Inches(2.3) + Inches(1.3) * i, CONTENT_W, Inches(1.15),
              title, body, color)

footer(s, 4)
add_notes(s, 'Selfieというブランド名を紹介。"自走"がキーワード。日本から指示し続ける状態は成功ではない。（3分）')


# ====================
# Slide 05: レンガ職人の話
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_WHITE)
logo_badge(s)
title_bar(s, "レンガ職人の話", "The Bricklayer's Story")
key_message(s, "同じ仕事でも「目的」が違うと成果が変わる", y=Inches(1.15))

bricks = [
    ("A", "ただ作業をしている", "（労働）", "Just working", C_GRAY, C_GRAY_LIGHT),
    ("B", "壁を作っている", "（職業・対価）", "Building a wall", C_BLUE, C_BLUE_LIGHT),
    ("C", "街（未来）を作っている", "（使命・Vision）", "Building a city", C_RED, C_RED_LIGHT),
]
card_w = Inches(3.5)
for i, (letter, text, sub, en, color, bg_color) in enumerate(bricks):
    x = ML + Inches(3.9) * i
    y_pos = Inches(2.3)
    # カード
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, card_w, Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if i == 2:  # C だけ赤枠
        card.line.color.rgb = C_RED
        card.line.width = Pt(2)
    else:
        card.line.color.rgb = C_BORDER
        card.line.width = Pt(0.5)
    # ラベル
    add_text(s, x + Inches(0.2), y_pos + Inches(0.2), card_w - Inches(0.4), Inches(0.4),
             f"Renovation {letter}", size=11, color=color, font=FONT_EN)
    add_text(s, x + Inches(0.2), y_pos + Inches(0.7), card_w - Inches(0.4), Inches(0.5),
             text, size=15, bold=True, color=C_TITLE, font=FONT)
    add_text(s, x + Inches(0.2), y_pos + Inches(1.3), card_w - Inches(0.4), Inches(0.3),
             sub, size=10, color=C_SUB, font=FONT)
    add_text(s, x + Inches(0.2), y_pos + Inches(1.7), card_w - Inches(0.4), Inches(0.3),
             en, size=9, color=C_GRAY, font=FONT_EN)
    # 矢印（最後以外）
    if i < 2:
        add_text(s, x + card_w + Inches(0.05), y_pos + Inches(0.8), Inches(0.4), Inches(0.5),
                 "▶", size=18, color=C_SUB, align=PP_ALIGN.CENTER, font=FONT)

bottom_bar(s, "私たちは「未来の仕組み」を作る側でいたい。",
           "We aspire to be the creators of future systems.",
           y=Inches(5.4))
footer(s, 5)
add_notes(s, "A/Bを否定しない。目的意識で成果が変わるメッセージに絞る。「みなさんは今どれに近い？」と軽く問いかけ。（2分）")


# ====================
# Slide 06: なぜMVVが必要か
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_WHITE)
logo_badge(s)
title_bar(s, "なぜMVVが必要か", "Why MVV Matters")
key_message(s, "正解がすぐに見えないからこそ、\n判断軸が必要", y=Inches(1.15))

reasons = [
    ("文化・常識・言語が違うほど、意思決定が割れる",
     "前提共有に時間がかかり、現場が止まる原因になる", C_NAVY, False),
    ("迷いはスピードを殺し、信頼を削る",
     "遅い判断は機会損失だけでなく、パートナーの不信感を生む", C_BLUE, False),
    ("MVVは「迷ったときに戻る場所」",
     "個人の感覚ではなく、組織としての共通の判断基準を持つ\nMVV is the place you return to when you're lost.", C_RED, True),
]
for i, (title, desc, color, highlight) in enumerate(reasons):
    y_pos = Inches(2.3) + Inches(1.35) * i
    card = info_card(s, ML, y_pos, CONTENT_W, Inches(1.2),
                     title, desc, color)
    if highlight:
        card.line.color.rgb = C_RED
        card.line.width = Pt(1.5)

footer(s, 6)
add_notes(s, "海外事業の文化の違いを具体例で。「MVVは飾りではなく実務ツール」と強調。（3分）")


# ====================
# Slide 07: WHAT WE AIM FOR
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_NAVY)
# 装飾サークル
circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Inches(-1), Inches(-1), Inches(8), Inches(8))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x22, 0x48, 0x6E)
circle.line.fill.background()
# メインタイトル
add_text(s, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
         "WHAT WE\nAIM FOR", size=56, bold=True, color=C_WHITE, font=FONT_EN,
         align=PP_ALIGN.CENTER)
# 赤ライン
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                           Inches(5.8), Inches(3.3), Inches(1.8), Inches(0.07))
line.fill.solid()
line.fill.fore_color.rgb = C_RED
line.line.fill.background()
# サブ
add_text(s, Inches(1.5), Inches(3.6), Inches(10), Inches(0.5),
         "何 を 目 指 す か", size=18, color=RGBColor(0xCC, 0xDD, 0xEE),
         align=PP_ALIGN.CENTER, font=FONT)

# 3つの要素
mvv_items = [
    ("Mission", "使命"),
    ("Vision", "将来像"),
    ("Values", "行動指針"),
]
for i, (en, ja) in enumerate(mvv_items):
    x = Inches(2.8) + Inches(2.8) * i
    # サークル
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(4.5), Inches(1.0), Inches(1.0))
    c.fill.background()
    c.line.color.rgb = C_WHITE
    c.line.width = Pt(2)
    # テキスト
    add_text(s, x - Inches(0.5), Inches(5.6), Inches(2.0), Inches(0.4),
             en, size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(s, x - Inches(0.5), Inches(6.0), Inches(2.0), Inches(0.3),
             ja, size=11, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER, font=FONT)
    # 矢印
    if i < 2:
        add_text(s, x + Inches(1.3), Inches(4.8), Inches(0.5), Inches(0.5),
                 "›", size=24, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER, font=FONT_EN)

add_notes(s, "セクション見出し。長くしない。「ここから具体的なMVVの中身に入ります」と予告。（1分）")


# ====================
# Slide 08: Mission
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_WHITE)
logo_badge(s)
title_bar(s, "Mission（使命）")

# 引用符
add_text(s, Inches(10.5), Inches(1.2), Inches(1), Inches(0.8),
         "❝", size=48, color=RGBColor(0xDD, 0xDD, 0xEE), font=FONT)
# 青い左バー
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          ML, Inches(1.5), Inches(0.08), Inches(2.0))
bar.fill.solid()
bar.fill.fore_color.rgb = C_NAVY
bar.line.fill.background()
# Mission本文
add_text(s, ML + Inches(0.25), Inches(1.5), Inches(10), Inches(0.8),
         "笑顔と可能性がめぐる未来をつくる。", size=30, bold=True, color=C_TITLE, font=FONT)
add_text(s, ML + Inches(0.25), Inches(2.3), Inches(10), Inches(0.5),
         "Create a future where smiles and possibilities grow.",
         size=14, color=C_SUB, font=FONT_EN)
# 説明
add_text(s, ML + Inches(0.25), Inches(3.0), Inches(10), Inches(0.4),
         "洗濯機を届けるのではなく、笑顔と可能性が人・地域・未来へ広がっていく状態をつくる。",
         size=12, color=C_BODY, font=FONT)

# 3つの補足カード
supplements = [
    ("「笑顔」", "生活が前に進む実感", "Smiles = the feeling of\nlife moving forward"),
    ("「可能性」", "時間が戻り、選択肢が増えること", "Possibilities = more time,\nmore choices"),
    ("「めぐる」", "一人で終わらず、次へつながること", "Grow = spreading to the\nnext person, next region"),
]
for i, (label, desc, en) in enumerate(supplements):
    x = ML + Inches(3.95) * i
    info_card(s, x, Inches(3.8), Inches(3.7), Inches(2.0),
              label, desc, C_NAVY, en)

footer(s, 8)
add_notes(s, "Missionを声に出して読む。参加者にも。「笑顔」「可能性」「めぐる」の3語を分解。「洗濯機を届けることがゴールじゃない」。英語版も読み上げる。（3分）")


# ====================
# Slide 09: Vision
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_WHITE)
logo_badge(s)
title_bar(s, "Vision（将来像）")

# 引用符
add_text(s, Inches(10.5), Inches(1.2), Inches(1), Inches(0.8),
         "❝", size=48, color=RGBColor(0xDD, 0xDD, 0xEE), font=FONT)
# 青い左バー
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          ML, Inches(1.5), Inches(0.08), Inches(2.0))
bar.fill.solid()
bar.fill.fore_color.rgb = C_NAVY
bar.line.fill.background()
# Vision本文
add_text(s, ML + Inches(0.25), Inches(1.5), Inches(10), Inches(0.8),
         "誰もが、自分の時間を力に変え、\n前へ進める世界をつくる。",
         size=26, bold=True, color=C_TITLE, font=FONT)
add_text(s, ML + Inches(0.25), Inches(2.6), Inches(10), Inches(0.5),
         "Build a world where people turn time into strength and move forward.",
         size=13, color=C_SUB, font=FONT_EN)
add_text(s, ML + Inches(0.25), Inches(3.1), Inches(10), Inches(0.4),
         "時間はただ減らす対象ではなく、人を前に進める力そのもの。",
         size=11, color=C_BODY, font=FONT)

# 3つの情景カード
scenes = [
    ("労働からの解放", '洗濯が"丸一日仕事"\nではなくなる', "Freedom from labor", C_NAVY),
    ("未来への投資", "子どもは学びへ\n大人は仕事や挑戦へ", "Investment in the future", C_BLUE),
    ("豊かさの実感", "家族の時間が増え\n本来の生活が戻ってくる", "Sense of abundance", C_RED),
]
for i, (title, desc, en, color) in enumerate(scenes):
    x = ML + Inches(3.95) * i
    info_card(s, x, Inches(3.6), Inches(3.7), Inches(1.7),
              title, desc, color, en)

bottom_bar(s, "「当たり前」が変わる ＝ 社会が変わる",
           "Changing the norm means changing society.", y=Inches(5.8))
footer(s, 9)
add_notes(s, "Visionを声に出して読む。「時間を力に変える」の具体例を情景として語る。英語版も読み上げる。（3分）")


# ====================
# Slide 10: ブランド思想
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_NAVY_DARK)
# メインメッセージ
add_text(s, Inches(1.0), Inches(1.5), Inches(11), Inches(1.5),
         "Time is not money —\nit's the power to\nchange the world.",
         size=42, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, font=FONT_EN)
# 赤ライン
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                           Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.06))
line.fill.solid()
line.fill.fore_color.rgb = C_RED
line.line.fill.background()
# サブ
add_text(s, Inches(1.0), Inches(3.8), Inches(11), Inches(0.5),
         "支援じゃない。一緒に、事業をつくる。",
         size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, font=FONT)
add_text(s, Inches(1.0), Inches(4.3), Inches(11), Inches(0.4),
         "Not aid. Building a business together. That's Life Breeze.",
         size=12, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER, font=FONT_EN)
# ファクト
add_text(s, Inches(2.5), Inches(5.2), Inches(8), Inches(0.3),
         "手洗い洗濯に週約8時間かかる現実を変える",
         size=11, color=RGBColor(0x99, 0xAA, 0xCC), align=PP_ALIGN.CENTER, font=FONT)
add_text(s, Inches(2.5), Inches(5.5), Inches(8), Inches(0.3),
         "戻った時間が、教育・仕事・家族・夢につながる",
         size=11, color=RGBColor(0x99, 0xAA, 0xCC), align=PP_ALIGN.CENTER, font=FONT)
add_notes(s, "英語のまま読み上げ、「時間はお金ではない。世界を変える力だ」と訳す。Life Breezeは支援団体ではなく事業会社。（3分）")


# ====================
# Slide 11: Values概要
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_WHITE)
logo_badge(s)
title_bar(s, "Values（行動指針）", "Our Values")
key_message(s, "5つの行動指針が、私たちの動き方を決める", y=Inches(1.15))

values = [
    ("1", "まず現場に立つ。", "Start from the ground."),
    ("2", "知ろうとし続ける。", "Stay curious. Stay humble."),
    ("3", "笑顔につながる価値をつくる。", "Create value that leads to smiles."),
    ("4", "受けた価値を、次へつなぐ。", "Pay forward the value you receive."),
    ("5", "利他で考え、続く形にする。", "Think beyond yourself. Make it last."),
]
colors = [C_NAVY, C_BLUE, C_RED, C_GREEN, C_AMBER]
for i, (num, ja, en) in enumerate(values):
    y_pos = Inches(2.1) + Inches(0.95) * i
    step_card(s, ML, y_pos, CONTENT_W, int(num), ja, "", colors[i], en)

footer(s, 11)
add_notes(s, "5つを一覧で見せ全体像を把握。個別の説明はこの後。（2分）")


# ====================
# Slide 12: Value 1-3
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_WHITE)
logo_badge(s)
title_bar(s, "Values — 現場・好奇心・笑顔")

vals = [
    ("Value 1: まず現場に立つ。", "Start from the ground.",
     "アフリカの現実は、日本の会議室では見えない。\nだから私たちは、まず現地に足を運ぶ。",
     "インターンの実践: 情報は二次情報で止めず、自分の目で確かめる。", C_NAVY),
    ("Value 2: 知ろうとし続ける。", "Stay curious. Stay humble.",
     "自分たちはまだ知らないことだらけだ。\nだからこそ、学び続けることをやめない。",
     "インターンの実践: 「わかったつもり」にならず、会議で最低1回は質問する。", C_BLUE),
    ("Value 3: 笑顔につながる価値をつくる。", "Create value that leads to smiles.",
     "作業を終わらせることが仕事ではない。\nその先に誰かの笑顔があるかどうかを問い続ける。",
     "インターンの実践: タスクの先にいる「誰か」を常に意識する。", C_RED),
]
for i, (title, en, desc, practice, color) in enumerate(vals):
    y_pos = Inches(1.2) + Inches(2.0) * i
    # カード
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               ML, y_pos, CONTENT_W, Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = C_WHITE
    card.line.color.rgb = C_BORDER
    card.line.width = Pt(0.5)
    # 左バー
    side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               ML, y_pos, Inches(0.06), Inches(1.8))
    side.fill.solid()
    side.fill.fore_color.rgb = color
    side.line.fill.background()
    # テキスト
    add_text(s, ML + Inches(0.2), y_pos + Inches(0.1), Inches(6), Inches(0.3),
             title, size=13, bold=True, color=C_TITLE, font=FONT)
    add_text(s, ML + Inches(0.2), y_pos + Inches(0.38), Inches(6), Inches(0.25),
             en, size=10, color=C_SUB, font=FONT_EN)
    add_text(s, ML + Inches(0.2), y_pos + Inches(0.65), Inches(6), Inches(0.5),
             desc, size=10, color=C_BODY, font=FONT)
    # 実践カード
    prac_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(7.5), y_pos + Inches(0.15), Inches(5.0), Inches(1.5))
    prac_bg.fill.solid()
    prac_bg.fill.fore_color.rgb = C_NAVY_LIGHT
    prac_bg.line.fill.background()
    add_text(s, Inches(7.7), y_pos + Inches(0.25), Inches(4.6), Inches(1.3),
             practice, size=10, color=C_NAVY, font=FONT)

footer(s, 12)
add_notes(s, "各Value 1分で説明。「インターンの実践」を特に強調。（3分）")


# ====================
# Slide 13: Value 4-5
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_WHITE)
logo_badge(s)
title_bar(s, "Values — つなぐ・利他")

vals2 = [
    ("Value 4: 受けた価値を、次へつなぐ。", "Pay forward the value you receive.",
     "学び・機会・成功体験を自分のところで止めない。\n次の人へ、次の地域へ、次の国へ渡していく。",
     "インターンの実践: 学んだことを言語化し、Notionやチャットでチームに共有する。", C_GREEN),
    ("Value 5: 利他で考え、続く形にする。", "Think beyond yourself. Make it last.",
     "自分だけが得をする仕組みは長続きしない。\n現地・地域・未来にとっても良い形を、諦めずに設計し続ける。",
     "インターンの実践: 「自分がいなくなっても続くか？」を判断基準にする。引き継ぎ資料を残す。", C_AMBER),
]
for i, (title, en, desc, practice, color) in enumerate(vals2):
    y_pos = Inches(1.2) + Inches(2.8) * i
    # カード
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               ML, y_pos, CONTENT_W, Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = C_WHITE
    card.line.color.rgb = C_BORDER
    card.line.width = Pt(0.5)
    # 左バー
    side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               ML, y_pos, Inches(0.06), Inches(2.5))
    side.fill.solid()
    side.fill.fore_color.rgb = color
    side.line.fill.background()
    # テキスト
    add_text(s, ML + Inches(0.2), y_pos + Inches(0.15), Inches(6), Inches(0.3),
             title, size=14, bold=True, color=C_TITLE, font=FONT)
    add_text(s, ML + Inches(0.2), y_pos + Inches(0.48), Inches(6), Inches(0.25),
             en, size=11, color=C_SUB, font=FONT_EN)
    add_text(s, ML + Inches(0.2), y_pos + Inches(0.85), Inches(6), Inches(0.7),
             desc, size=11, color=C_BODY, font=FONT)
    # 実践カード
    prac_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(7.5), y_pos + Inches(0.3), Inches(5.0), Inches(1.8))
    prac_bg.fill.solid()
    prac_bg.fill.fore_color.rgb = C_NAVY_LIGHT
    prac_bg.line.fill.background()
    add_text(s, Inches(7.7), y_pos + Inches(0.5), Inches(4.6), Inches(1.5),
             practice, size=11, color=C_NAVY, font=FONT)

footer(s, 13)
add_notes(s, "Value 4は「学びを独占しない」。Value 5は「自分がいなくなっても続くか？」。自分中心ではなくつながりの中で考える。（3分）")


# ====================
# Slide 14: 迷ったらこの順番
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_WHITE)
logo_badge(s)
title_bar(s, "迷ったらこの順番", "Decision Framework")

steps_data = [
    ("01", "Mission", "最優先 / Top Priority",
     "Missionに沿うか？（笑顔と可能性がめぐる未来につながるか）",
     "Does it align with our Mission?", C_NAVY),
    ("02", "安全 / 法令", "絶対遵守 / Non-negotiable",
     "安全・法令・倫理はOKか？（誰かを傷つけないか）",
     "Is it safe, legal, and ethical?", C_RED),
    ("03", "信頼 / Trust", "",
     "信頼を積むか？（約束・透明性・説明責任）",
     "Does it build trust?", C_GREEN),
    ("04", "再現性 / Reproducibility", "",
     "再現できるか？（型として残せるか？）",
     "Can it be replicated?", C_GRAY),
]
for i, (num, label, badge, desc_ja, desc_en, color) in enumerate(steps_data):
    y_pos = Inches(1.3) + Inches(1.35) * i
    # カード
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               ML, y_pos, CONTENT_W, Inches(1.2))
    card.fill.solid()
    card.fill.fore_color.rgb = C_WHITE
    card.line.color.rgb = C_BORDER
    card.line.width = Pt(0.5)
    # 左色帯
    side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               ML, y_pos, Inches(0.06), Inches(1.2))
    side.fill.solid()
    side.fill.fore_color.rgb = color
    side.line.fill.background()
    # STEPバッジ
    sbadge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 ML + Inches(0.2), y_pos + Inches(0.15),
                                 Inches(0.7), Inches(0.35))
    sbadge.fill.solid()
    sbadge.fill.fore_color.rgb = C_GRAY_LIGHT
    sbadge.line.fill.background()
    tf_sb = sbadge.text_frame
    tf_sb.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_sb = tf_sb.paragraphs[0]
    p_sb.text = f"STEP"
    p_sb.font.size = Pt(8)
    p_sb.font.color.rgb = C_SUB
    p_sb.font.name = FONT_EN
    p_sb.alignment = PP_ALIGN.CENTER
    # 番号
    num_shape = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    ML + Inches(1.1), y_pos + Inches(0.1),
                                    Inches(0.45), Inches(0.45))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = color
    num_shape.line.fill.background()
    tf_n = num_shape.text_frame
    tf_n.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_n = tf_n.paragraphs[0]
    p_n.text = num
    p_n.font.size = Pt(14)
    p_n.font.bold = True
    p_n.font.color.rgb = C_WHITE
    p_n.font.name = FONT_EN
    p_n.alignment = PP_ALIGN.CENTER
    # ラベル
    add_text(s, ML + Inches(1.8), y_pos + Inches(0.1), Inches(3), Inches(0.35),
             label, size=14, bold=True, color=color, font=FONT)
    if badge:
        add_text(s, ML + Inches(4.5), y_pos + Inches(0.15), Inches(2.5), Inches(0.25),
                 badge, size=8, color=C_SUB, font=FONT_EN)
    # 説明
    add_text(s, ML + Inches(1.8), y_pos + Inches(0.5), Inches(6), Inches(0.3),
             desc_ja, size=10, color=C_BODY, font=FONT)
    add_text(s, ML + Inches(1.8), y_pos + Inches(0.8), Inches(6), Inches(0.25),
             desc_en, size=9, color=C_GRAY, font=FONT_EN)
    # 矢印
    if i < 3:
        add_text(s, Inches(6.2), y_pos + Inches(1.15), Inches(0.5), Inches(0.3),
                 "↓", size=14, color=C_SUB, align=PP_ALIGN.CENTER, font=FONT)

footer(s, 14)
add_notes(s, "4ステップを順に説明。「暗記しなくてOK。迷ったときにこのスライドを見返せばいい。」（3分）")


# ====================
# Slide 15: Workshop 01
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_WHITE)
logo_badge(s)
workshop_header(s, 1, "一番刺さったValueを選ぶ", C_NAVY, 5)

# メインメッセージ
add_text(s, Inches(2.0), Inches(2.0), Inches(9), Inches(0.8),
         "一番刺さったValueを1つ選び、\n理由を書いてください",
         size=22, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER, font=FONT)

# Q1
q1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(1.5), Inches(3.2), Inches(10), Inches(1.3))
q1.fill.solid()
q1.fill.fore_color.rgb = C_GRAY_LIGHT
q1.line.fill.background()
add_text(s, Inches(1.8), Inches(3.3), Inches(1.0), Inches(0.35),
         "Q1", size=14, bold=True, color=C_NAVY, font=FONT_EN)
add_text(s, Inches(2.8), Inches(3.3), Inches(7), Inches(0.35),
         "どのValue？ / Which Value?", size=13, bold=True, color=C_TITLE, font=FONT)
add_text(s, Inches(2.8), Inches(3.7), Inches(7), Inches(0.5),
         "5つのValueから、最も共感したもの・重要だと思ったものを1つ選ぶ",
         size=10, color=C_BODY, font=FONT)
# Q2
q2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(1.5), Inches(4.8), Inches(10), Inches(1.3))
q2.fill.solid()
q2.fill.fore_color.rgb = C_GRAY_LIGHT
q2.line.fill.background()
add_text(s, Inches(1.8), Inches(4.9), Inches(1.0), Inches(0.35),
         "Q2", size=14, bold=True, color=C_NAVY, font=FONT_EN)
add_text(s, Inches(2.8), Inches(4.9), Inches(7), Inches(0.35),
         "なぜ刺さった？ / Why?", size=13, bold=True, color=C_TITLE, font=FONT)
add_text(s, Inches(2.8), Inches(5.3), Inches(7), Inches(0.5),
         "理由を言語化する（例: 過去の経験から / 今の課題だから / 強みにしたいから）",
         size=10, color=C_BODY, font=FONT)

footer(s, 15)
add_notes(s, "「正解はありません。直感でOK。」タイマーを見せる。3分で書いて2分で共有。（5分）")


# ====================
# Slide 16: Workshop 02
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_WHITE)
logo_badge(s)
workshop_header(s, 2, "ケーススタディ", C_GREEN, 5)

# ケース
case_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.5))
case_bg.fill.solid()
case_bg.fill.fore_color.rgb = C_GRAY_LIGHT
case_bg.line.fill.background()
add_text(s, Inches(0.8), Inches(1.7), Inches(0.8), Inches(0.35),
         "?", size=20, bold=True, color=C_GREEN, font=FONT_EN)
add_text(s, Inches(1.5), Inches(1.65), Inches(2), Inches(0.3),
         "CASE STUDY", size=9, bold=True, color=C_GREEN, font=FONT_EN)
add_text(s, Inches(1.5), Inches(1.95), Inches(10), Inches(0.9),
         "あなたはインターンとして資料を作成中。締切は明日。\n上司から「今の資料でいいから今すぐ送って」と言われた。\nでも、数字の確認がまだ終わっていない。あなたならどうする？",
         size=11, color=C_BODY, font=FONT)

# 4ステップ
fw_steps = [
    ("01", "Mission", "笑顔と可能性につながる？", C_NAVY),
    ("02", "安全/法令", "誤った数字を出すリスクは？", C_RED),
    ("03", "信頼", "相手の信頼を積める？", C_GREEN),
    ("04", "再現性", "この対応は型にできる？", C_GRAY),
]
for i, (num, label, desc, color) in enumerate(fw_steps):
    x = Inches(0.5) + Inches(3.1) * i
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               x, Inches(3.4), Inches(2.9), Inches(1.2))
    card.fill.solid()
    card.fill.fore_color.rgb = C_WHITE
    card.line.color.rgb = color
    card.line.width = Pt(1)
    add_text(s, x + Inches(0.15), Inches(3.5), Inches(0.4), Inches(0.3),
             num, size=12, bold=True, color=color, font=FONT_EN)
    add_text(s, x + Inches(0.55), Inches(3.5), Inches(2), Inches(0.3),
             label, size=12, bold=True, color=color, font=FONT)
    add_text(s, x + Inches(0.15), Inches(3.9), Inches(2.6), Inches(0.5),
             desc, size=10, color=C_BODY, font=FONT)

# 結論
conclusion = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(1.0), Inches(4.9), Inches(11.3), Inches(1.5))
conclusion.fill.background()
conclusion.line.color.rgb = C_NAVY
conclusion.line.width = Pt(1.5)
conclusion.line.dash_style = 4  # dash
add_text(s, Inches(4.5), Inches(5.0), Inches(4), Inches(0.3),
         "YOUR DECISION", size=10, bold=True, color=C_RED, align=PP_ALIGN.CENTER, font=FONT_EN)
add_text(s, Inches(1.5), Inches(5.3), Inches(10), Inches(0.4),
         "結論：条件付きでどう進めるか  (Yes if / No because + option)",
         size=12, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER, font=FONT)

footer(s, 16)
add_notes(s, "インターンの日常に近いケース。正解は1つではない。「Yes, if...」か「No, because... & option」の型で。2分で考え、2分で共有。（5分）")


# ====================
# Slide 17: Workshop 03
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_WHITE)
logo_badge(s)
workshop_header(s, 3, "30日アクション宣言", C_RED, 7)

# メイン
add_text(s, Inches(2.0), Inches(1.8), Inches(9), Inches(0.8),
         '30日でやる"たった1つ"の行動を宣言する',
         size=22, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER, font=FONT)
add_text(s, Inches(2.0), Inches(2.5), Inches(9), Inches(0.3),
         "小さくても「完了」できるサイズに切り出すことが重要です",
         size=11, color=C_SUB, align=PP_ALIGN.CENTER, font=FONT)

# 例カード
examples = [
    ("リサーチ", "Life Breezeの展開国に\nついて1つ調べ、\nチームに共有する", C_NAVY),
    ("学びの共有", "毎週の学びを1つ、\nテキストでチームに\n送る", C_GREEN),
    ("現場理解", "現地スタッフの日報を\n1週間読み、\n気づきをまとめる", C_AMBER),
]
for i, (label, desc, color) in enumerate(examples):
    x = Inches(1.5) + Inches(3.8) * i
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               x, Inches(3.2), Inches(3.3), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = C_WHITE
    card.line.color.rgb = color
    card.line.width = Pt(1)
    # ラベルバッジ
    lbl = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              x + Inches(0.15), Inches(3.35), Inches(1.5), Inches(0.3))
    lbl.fill.solid()
    lbl.fill.fore_color.rgb = color
    lbl.line.fill.background()
    tf_l = lbl.text_frame
    tf_l.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_l = tf_l.paragraphs[0]
    p_l.text = f"  例：{label}"
    p_l.font.size = Pt(9)
    p_l.font.bold = True
    p_l.font.color.rgb = C_WHITE
    p_l.font.name = FONT
    # 内容
    add_text(s, x + Inches(0.2), Inches(3.8), Inches(2.9), Inches(1.5),
             desc, size=12, bold=True, color=C_TITLE, font=FONT)

# ボトムバー
bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.7))
bar.fill.solid()
bar.fill.fore_color.rgb = C_GRAY_LIGHT
bar.line.fill.background()
add_text(s, Inches(2.0), Inches(6.1), Inches(9), Inches(0.4),
         "私は30日以内に「_______________」をやります。それは Value ___につながる行動です。",
         size=11, color=C_BODY, font=FONT)

footer(s, 17)
add_notes(s, "「大きなことでなくていい。小さくても完了できるもの。」タイマーで4分。抽象的な宣言には「具体的にすると？」と質問。2-3名に発表。（7分）")


# ====================
# Slide 18: まとめ
# ====================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, C_WHITE)
logo_badge(s)
title_bar(s, "まとめ", "Summary")
# メイン
add_text(s, ML, Inches(1.2), CONTENT_W, Inches(0.7),
         '私たちは"仕組み"で\n笑顔と可能性がめぐる未来をつくる',
         size=20, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER, font=FONT)

# 3カラムカード
summary = [
    ("Mission", "使命", "笑顔と可能性が\nめぐる未来をつくる。",
     "Create a future where\nsmiles and possibilities grow.", C_NAVY),
    ("Vision", "将来像", "誰もが、自分の時間を\n力に変え、前へ進める\n世界をつくる。",
     "Build a world where people\nturn time into strength.", C_BLUE),
    ("Values", "行動指針", "現場 / 好奇心 / 笑顔\nつなぐ / 利他",
     "5 values that define\nhow we act.", C_RED),
]
for i, (en, ja, desc, desc_en, color) in enumerate(summary):
    x = ML + Inches(3.95) * i
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               x, Inches(2.2), Inches(3.7), Inches(3.2))
    card.fill.solid()
    card.fill.fore_color.rgb = C_WHITE
    card.line.color.rgb = C_BORDER
    card.line.width = Pt(0.5)
    # アイコンサークル
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 x + Inches(1.45), Inches(2.4), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = C_NAVY_LIGHT
    circle.line.fill.background()
    # ラベル
    add_text(s, x, Inches(3.3), Inches(3.7), Inches(0.35),
             en, size=14, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(s, x, Inches(3.6), Inches(3.7), Inches(0.25),
             ja, size=10, color=C_SUB, align=PP_ALIGN.CENTER, font=FONT)
    # 説明
    add_text(s, x + Inches(0.2), Inches(3.95), Inches(3.3), Inches(0.8),
             desc, size=11, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER, font=FONT)
    add_text(s, x + Inches(0.2), Inches(4.7), Inches(3.3), Inches(0.5),
             desc_en, size=9, color=C_SUB, align=PP_ALIGN.CENTER, font=FONT_EN)

# NEXT ACTION バー
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          Inches(0), Inches(5.8), SLIDE_W, Inches(1.2))
bar.fill.solid()
bar.fill.fore_color.rgb = C_NAVY_DARK
bar.line.fill.background()
add_text(s, Inches(1.0), Inches(5.85), Inches(2.5), Inches(0.35),
         "NEXT ACTION", size=10, bold=True, color=C_RED, font=FONT_EN)
add_text(s, Inches(1.0), Inches(6.15), Inches(10), Inches(0.5),
         "明日から、30日アクションを1回実行する。",
         size=20, bold=True, color=C_WHITE, font=FONT)
add_text(s, Inches(1.0), Inches(6.55), Inches(10), Inches(0.3),
         "Starting tomorrow — execute your 30-day action once.",
         size=11, color=RGBColor(0xAA, 0xBB, 0xDD), font=FONT_EN)

footer(s, 18)
add_notes(s, "Mission→Vision→Valuesを一言ずつ振り返る。「MVVは飾りではなく、明日から使う判断基準です。30日後に振り返りましょう。」（5分）")


# ====================
# 保存
# ====================
out_dir = os.path.join(os.path.dirname(__file__), "slides")
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "Life_Breeze_MVV_Training_Intern_v0_1.pptx")
prs.save(out_path)
print(f"[OK] Generated: {out_path}")
print(f"     Slides: {len(prs.slides)}")
